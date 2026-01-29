#!/usr/bin/env python3
"""Generate a PPTX deck from a folder of slide images.

Usage:
  python3 tools/generate_pptx_from_images.py \
    --input-dir deck-assets \
    --pattern 'slide-*.jpg' \
    --output tmp/autorenta_deck.pptx

Notes:
- Produces a 16:9 deck (13.333 x 7.5 inches).
- Each matching image becomes a full-bleed slide.
"""

from __future__ import annotations

import argparse
import glob
import os
import sys
from pathlib import Path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate PPTX from slide images")
    parser.add_argument("--input-dir", default="deck-assets", help="Directory containing slide images")
    parser.add_argument("--pattern", default="slide-*.jpg", help="Glob pattern (relative to input-dir)")
    parser.add_argument("--output", default="tmp/autorenta_deck.pptx", help="Output PPTX path")
    parser.add_argument(
        "--sort",
        default="name",
        choices=["name"],
        help="Sorting method for images (currently: name)",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore
    except Exception as exc:  # pragma: no cover
        print(
            "Missing dependency: python-pptx. Install with:\n"
            "  python3 -m pip install --user python-pptx\n",
            file=sys.stderr,
        )
        print(f"Import error: {exc}", file=sys.stderr)
        return 2

    input_dir = Path(args.input_dir)
    output_path = Path(args.output)
    pattern = str(input_dir / args.pattern)

    images = sorted(Path(p).resolve() for p in glob.glob(pattern))
    if not images:
        print(f"No images found for pattern: {pattern}", file=sys.stderr)
        return 1

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Force widescreen 16:9 (PowerPoint default widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for image_path in images:
        slide = prs.slides.add_slide(blank_layout)
        # Full-bleed. Assumes input images are already 16:9.
        slide.shapes.add_picture(str(image_path), 0, 0, width=slide_w, height=slide_h)

    prs.save(str(output_path))
    print(f"Wrote {len(images)} slides to: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
